from pptx import Presentation
from PIL import Image
# import sc????
import pandas as pd
from subprocess import call
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import requests
import json
from datetime import datetime as dt

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches,Cm, Pt

import get_bus_info

from pptx.enum.chart import XL_TICK_MARK

from pptx.enum.chart import XL_CATEGORY_TYPE

from financial_info import  *

import scraperbot
from summa import summarizer

from get_company_logo import get_company_logo

class Comp_prof():

	def __init__(self,ticker,company,file = '/Users/adityakaushik/ppt_project_master/sample_MA.pptx'):



		# company is name
		# ticker is the yahoo finance ticker or is nothing if the company is private
		# rtc is the reuters code
		currency_dict = {
			'NYSE': '$',
			'L': '£',
			'PA': '€',
			'DE': '€',
			'MI': '€'
		}


		self.presentation = Presentation(file)
		self.company = company
		#try:
		if ticker == '':
			self.rtc = 'Private'
			self.ticker = None
			self.currency = ''
			print('setting ticker to None')
		elif not '.' in ticker:
			# american company
			self.ticker = ticker

			self.rtc = ticker
			self.currency = '$'
		else:
			# otherwise the ticker is non-American
			self.ticker = ticker
			try:
				# try to find the currency
				exchange = ticker.split('.')[1].upper()
				try:
					self.currency = currency_dict[exchange]
				except:
					self.currency = 'INSERT CURRENCY HERE'
				if exchange == 'L':
					# then the comapny is british
					self.rtc = ticker
				else:
					try:
						print('going to search for company')
						self.rtc = scraperbot.search_company(company)
						print('rtc code is '+ self.rtc)
					except:
						print('searching company failed')
						self.rtc = 'Private'
			except:
				self.rtc = 'Private'
				self.currency = 'INSERT CURRENCY HERE'

		#except:
		#	print('exception at initialization')
		#	# if we can't find the RTC code then the company is clearly private
		#	self.ticker = ticker
		#	self.rtc = 'Private'

		# need to come up with a way to get the currency



# prof =comp_prof('sample_MA.pptx')


	def run(self,slide_num =0):
		# returns presentation object
		self.table_filler()
		self.insert_key_people()
		self.insert_business_description()
		self.insert_recent_news()
		self.insert_quote_table()
		self.insert_key_shareholders()
		try:
			self.add_logo()
		except:
			print('unable to retrieve logo')
		company = self.company

		presentation = self.presentation

		presentation = Comp_prof.set_company_name(presentation,company)
		presentation = Comp_prof.add_stock_chart(presentation,self.ticker)
		return presentation	



	def find_placeholder_order(profile):
		#  useful for figuring out which placeholder corresponds to which  
		# profile is a comp_prof slide  object
		count =0

		for shapes in profile.shapes:
			if not shapes.has_text_frame:
				count += 1
				continue
			par = shapes.text_frame.paragraphs[0]
			par.clear()
			count += 1
			par.text = str(count)









	def set_company_name(ppt,Company_name):
		# Company Name is a string which represents the company's name
		# ppt is the Powerpoint object given. This particular method is build towards the template I've used
		# in general, it's imperative to figure out which placeholder indices you want to change


		slides = ppt.slides[0]
		title = slides.shapes[0]

		try:
			title.text = Company_name
		except:
			print('something went wrong with title')

		return ppt






	def insert_key_people(self):
		# People is a list of tuples. each tuple contains the person's name and then a description
		if self.rtc == "Private":
			return []
		else:
			#try:
				people = scraperbot.get_officer_information(self.rtc)

				slides = self.presentation.slides[0]
				people_box = slides.shapes[5].text_frame

				# we use 8 because that's the particular text frame we want to go to



				par = people_box.paragraphs[0]
				par.font.size = Pt(8)
				for peeps in people:
					par.font.size = Pt(8)
					name = peeps[0]
					description = peeps[1]


					par.clear()

					run = par.add_run()
					run.font.size = Pt(8)


					run3 = par.add_run()
					run3.text = description
					run3.font.size = Pt(8)
					font = run.font
					font.bold = True
					font.size = Pt(8)


					run3 = par.add_run()
					run3.font.size == Pt(8)
					run.text = name + ' - '
					run.font.size = Pt(8)
					par.add_line_break()
					par = people_box.add_paragraph()
		#	except:
		#		print('Error at comp_prof. insert_key_People() . self.rtc not == Private and unable to access officer information')
		#		return []










	def insert_business_description(self, slide_num =0):
		# ppt is the powerpoint we're acting on
		#slide is the index of the slide we're acting on
		# description is a list of sentences that describe the business

		# business description is indexed by placehodler 7

		description = get_bus_info.business_description_generator(company = self.company, rtc_code =self.rtc,number_sentences= 5)

		slides = self.presentation.slides[slide_num]
	
		bus_descrip= slides.shapes[4].text_frame
	
		bus_descrip.clear()
	
	
		par = bus_descrip.paragraphs[0]


		if not description == None:
			par.font.size = Pt(8)
			for sentences in description:
				par.font.size = Pt(8)
				run = par.add_run()
				run.text = sentences
				run.font.size = Pt(8)
				par = bus_descrip.add_paragraph()

				par.clear()






	def insert_recent_news(self,slide_num =0):
		#ppt is the powerpoint we are editing
	# slide_num is the slide we want to target 
	# news is a list of news items
		news = get_bus_info.news_summary_generator(company_name = self.company, rtc_code = self.rtc)

		slides = self.presentation.slides[slide_num]
	
		news_box= slides.shapes[3].text_frame
	
		news_box.clear()
	
	
		par = news_box.paragraphs[0]


	
		for sentences in news:
			par.font.size = Pt(8)
			run = par.add_run()
			run.text = sentences
			run.font.size = Pt(8)
			par = news_box.add_paragraph()
	
			par.clear()





	def add_stock_chart(ppt, ticker, slide_num =0):
	# data is a series indexed by date -- dates are all datetime objects!
	# ticker is a yahoo ticker-- note this only works with american and british companies
		if ticker == None:
			print('This is a private company so there will be no stock price')
			return ppt
		else:
			try:
				data = get_historic_price(ticker = ticker)
				slide = ppt.slides[slide_num]

				chart_data = CategoryChartData()

				chart_data.categories = data.index.tolist()
				chart_data.add_series(str(ticker),tuple(data.values.tolist()))


				## need to find a way to push the axis labels out -- ie delay the labels starting till a week after the first set of data



				x, y, cx, cy = Inches(.403), Inches(5.33), Inches(2.22), Inches(1.3)
				chart = slide.shapes.add_chart(
					XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
				).chart

				chart.font.size = Pt(8)
				chart.font.name = 'Arial'
			except:
				print('We recieved the ticker, but were unable to get the historic price. Invalid ticker')



			#  chart.value_axis.minimum_scale = (min(data.values.tolist()))

		date_axis = chart.category_axis
		date_axis.tick_labels.number_format = 'mmm-yy'
		date_axis.major_unit = 90
		chart.has_legend = False
		chart.show_legend_key = False

		chart.has_title = False
		chart.series[0].smooth = True

		return ppt

		# 3.33 cm by 5.72 cm length x width of desired chart

		#1.04 cm left 13.76 cm top desired location on slide




	def table_filler(self):
		# financials is a pandas dataframe with the 4 columns of the trailing years financials and rows
		# that include Revenue, Gross Profit, Gross Margin, EBITDA, EBIT, Operating Margin, Net Income, Net Margin
		if not self.ticker == None:

			try:
				financials = get_financials(self.ticker)
				dates = financials.columns.tolist()
				categories = financials.index.tolist()

				## This part of the code is unique to the profile I'm using
				## there could be multiple tables in the default slide; thus it's important to know which table you're actually filling!
				## To Do: Deprecate below for loop and simply KNOW which shape number is the table
				tables = []
				for shapes in self.presentation.slides[0].shapes:
					if shapes.has_table and shapes.table:
						tables.append(shapes.table)
				table = tables[0]

				# setting Currency / units in top left corner
				currency = self.currency

				par = table.cell(0, 0).text_frame.paragraphs[0]
				run = par.add_run()
				run.text = "("+currency+"m)"
				run.font.name = "arial"
				run.font.size = Pt(7)
				run.font.italic = True
				par.alignment = PP_ALIGN.RIGHT
				table.cell(0, 0).vertical_anchor = MSO_ANCHOR.MIDDLE


				# setting dates in first row
				for i in range(1, 5):
					try:
						par = table.cell(0, i).text_frame.paragraphs[0]
						run = par.add_run()
						# retreiving the relevant date from the dates list
						run.text = dates[i-1]
						run.font.name = "arial"
						run.font.size = Pt(8)
						par.alignment = PP_ALIGN.RIGHT
						table.cell(0, i).vertical_anchor = MSO_ANCHOR.MIDDLE
					except:
						print('not able to add additional dates')
				# setting categories (ie Revenue, etc.) in 1st column
				for j in range(9):
					try:
						margin_cells = [3, 6, 8]
						if j == 0:
							continue
						elif j in margin_cells:
							par = table.cell(j, 0).text_frame.paragraphs[0]
							run = par.add_run()
							# obtaining the relevant category from the categories list
							print(categories[j-1])
							run.text = categories[j-1]
							run.font.name = 'arial'
							run.font.size = Pt(9)
							run.font.italic = True
							table.cell(j, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
						else:
							par = table.cell(j, 0).text_frame.paragraphs[0]
							run = par.add_run()
							run.text = categories[j-1]
							run.font.name = 'arial'
							run.font.size = Pt(9)
							run.font.bold = True
							table.cell(j, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
					except:
						print('There were too many financials')

				for j in range(1, 5):
					numbers = financials[dates[j-1]].values.tolist()
					for i in range(1, 9):
						par = table.cell(i, j).text_frame.paragraphs[0]
						run = par.add_run()
						run.text = numbers[i-1]
						run.font.name = 'arial'
						run.font.size = Pt(8)
						par.alignment = PP_ALIGN.RIGHT
						table.cell(i, j).vertical_anchor = MSO_ANCHOR.MIDDLE
			except:
				print('Ticker is not equal to None and still recieved an error. At Table Filler')

	def add_logo(self):
		# Left Center: 21.cm
		# top center : .37cm
		# height 1.8cm



		ppt = self.presentation
		slide = ppt.slides[0]
		get_company_logo(self.company)
		file = self.company+'.png'

		# initially everything is in CM
		frame_height = 1.8
		im = Image.open(file)
		width, height = im.size
		try:
			dpi = im.info['dpi'][0]
		except:
			dpi = 72
		width = 2.54 / dpi * width
		height = 2.54 / dpi * height
		height_ratio = frame_height / height
		img_height = frame_height
		img_width = width * height_ratio

		x = img_width/2

		top_location = .37
		left_location = 21.17 -x
		slide.shapes.add_picture(file, left=Cm(left_location), top=Cm(top_location), height=Cm(img_height))
		subprocess.call(['rm', '/Users/adityakaushik/ppt_project_master/'+self.company+'.png'])

	def insert_quote_table(self):
		ticker = self.ticker
		presentation = self.presentation
		currency = self.currency
		relevant_info = get_quote_table(ticker)
		slide = presentation.slides[0]
		quote_table = slide.shapes[2].text_frame

		par  = quote_table.paragraphs[0]

		for i in relevant_info:
			run = par.add_run()
			if i[0] == 'PE Ratio (TTM)':
				run.text = i[0] + '\t' + i[1]
			else:
				run.text = i[0] + '\t' + currency + i[1]
			run.font.size = Pt(8)
			run.font.name = 'arial'
			par = quote_table.add_paragraph()
			par.clear()

	def insert_key_shareholders(self):
		if not self.ticker == None:

			try:
				shareholders = get_key_shareholders(self.ticker)
				names = shareholders['Holder'].tolist()
				categories = shareholders['% Out'].tolist()

				## This part of the code is unique to the profile I'm using
				## there could be multiple tables in the default slide; thus it's important to know which table you're actually filling!
				## To Do: Deprecate below for loop and simply KNOW which shape number is the table
				tables = []
				for shapes in self.presentation.slides[0].shapes:
					if shapes.has_table and shapes.table:
						tables.append(shapes.table)
				table = tables[1]

				for i in range(8):
					par = table.cell(i, 0).text_frame.paragraphs[0]
					run = par.add_run()
					run.text = names[i]
					run.font.size = Pt(7)
					run.font.name = 'arial'
					par.alignment = PP_ALIGN.LEFT

					par = table.cell(i, 1).text_frame.paragraphs[0]
					run = par.add_run()
					run.text = categories[i]
					run.font.size = Pt(7)
					run.font.name = 'arial'
					par.alignment = PP_ALIGN.RIGHT
			except:
				print('unable to fill Key Shareholders table')



if __name__ == 'main':
	run()

"""
		# we need to figure out if we should adjust size based
		# on height or based on width
		# to do thi we check the aspect ratios
		if width_ratio > height_ratio:
			print('reached part 1)')
			# then height will be the bound when reshaping, so we adjust height
			# converting everything to inches by dividing by 2.54
			img_height = frame_height
			img_width = width * height_ratio
			y = img_height / 2
			x = img_width / 2
			top_location = 3.92 - y
			left_location = 3.92 - x

			slide.shapes.add_picture(file, left=Cm(left_location), top=Cm(top_location), height=Cm(img_height))


		else:
			print('reached part 2')
			# need to convert everythign to inches by dividing by 2.54
			img_width = frame_width
			img_height = height * width_ratio
			y = img_height / 2
			x = img_width / 2
			top_location = 3.92 - y
			left_location = 3.92 - x

			slide.shapes.add_picture(file, left=Cm(left_location), top=Cm(top_location), width=Cm(img_width))
"""